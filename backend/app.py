import os
import logging
from flask import Flask, send_file, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import google.generativeai as genai
import re
import requests
from io import BytesIO
import base64
import pandas as pd
from dotenv import load_dotenv
import subprocess
import tempfile
from pathlib import Path
from fpdf import FPDF  # pip install fpdf
from PIL import Image  # pip install Pillow
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.action import PP_ACTION
from PIL import Image, ImageDraw, ImageFont

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Initialize Flask app and enable CORS
app = Flask(__name__)
CORS(app)

# Configure Gemini API
genai_api_key = os.environ.get('GENAI_API_KEY')
if not genai_api_key:
    app.logger.error("GENAI_API_KEY is not set in the environment!")
    raise ValueError("GENAI_API_KEY is required")
genai.configure(api_key=genai_api_key)

# Enhanced theme options with gradients and shadows
THEMES = {
    "corporate": {
        "gradient_start": RGBColor(0, 51, 102),
        "gradient_end": RGBColor(173, 216, 230),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(255, 215, 0),  # Gold accent
        "font_name": "Arial",
        "shadow": True,
        "variants": {
            "professional": {
                "layout": "clean",
                "font_name": "Calibri",
                "bullet_style": "square"
            },
            "creative": {
                "layout": "asymmetric",
                "font_name": "Segoe UI",
                "bullet_style": "circle"
            },
            "minimal": {
                "layout": "centered",
                "font_name": "Arial",
                "bullet_style": "dash"
            }
        }
    },
    "creative": {
        "gradient_start": RGBColor(147, 112, 219),
        "gradient_end": RGBColor(255, 182, 193),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(240, 240, 240),
        "accent_color": RGBColor(0, 255, 127),  # Green accent
        "font_name": "Montserrat",
        "shadow": True,
        "variants": {
            "professional": {
                "layout": "structured",
                "font_name": "Verdana",
                "bullet_style": "arrow"
            },
            "creative": {
                "layout": "dynamic",
                "font_name": "Comic Sans MS",
                "bullet_style": "star"
            },
            "minimal": {
                "layout": "simple",
                "font_name": "Tahoma",
                "bullet_style": "dot"
            }
        }
    },
    "minimal": {
        "gradient_start": RGBColor(245, 245, 245),
        "gradient_end": RGBColor(255, 255, 255),
        "title_color": RGBColor(33, 33, 33),
        "text_color": RGBColor(66, 66, 66),
        "accent_color": RGBColor(0, 150, 199),  # Blue accent
        "font_name": "Helvetica",
        "shadow": False,
        "variants": {
            "professional": {
                "layout": "grid",
                "font_name": "Helvetica",
                "bullet_style": "none"
            },
            "creative": {
                "layout": "freeform",
                "font_name": "Gill Sans",
                "bullet_style": "checkmark"
            },
            "minimal": {
                "layout": "spacious",
                "font_name": "Arial Narrow",
                "bullet_style": "hyphen"
            }
        }
    },
    "bold": {
        "gradient_start": RGBColor(255, 69, 0),
        "gradient_end": RGBColor(255, 215, 0),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(255, 255, 255),
        "accent_color": RGBColor(0, 0, 0),  # Black accent
        "font_name": "Impact",
        "shadow": True,
        "variants": {
            "professional": {
                "layout": "strong",
                "font_name": "Arial Black",
                "bullet_style": "triangle"
            },
            "creative": {
                "layout": "expressive",
                "font_name": "Impact",
                "bullet_style": "diamond"
            },
            "minimal": {
                "layout": "bold-minimal",
                "font_name": "Franklin Gothic",
                "bullet_style": "arrow"
            }
        }
    }
}

# Chart type mapping
CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

# Presentation styling constants
TITLE_FONT_SIZE = Pt(36)
CONTENT_FONT_SIZE = Pt(20)

# Function to apply gradient background
def apply_gradient(slide, start_color, end_color):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = start_color
    fill.gradient_stops[1].color.rgb = end_color
    fill.gradient_angle = 45

# Function to add shadow to text
def add_text_shadow(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            add_shadow_to_run(run)

def add_shadow_to_run(run):
    rPr = run._r.get_or_add_rPr()
    effectLst = rPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = OxmlElement('a:effectLst')
        rPr.append(effectLst)
    
    outerShdw = effectLst.find(qn('a:outerShdw'))
    if outerShdw is None:
        outerShdw = OxmlElement('a:outerShdw')
        outerShdw.set('blurRad', '40000')
        outerShdw.set('dist', '20000')
        outerShdw.set('dir', '5400000')
        outerShdw.set('rotWithShape', '0')
        
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '323232')
        
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '50000')
        srgbClr.append(alpha)
        
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)

def process_titles(text):
    try:
        lines = text.strip().split('\n')
        titles = [line.strip() for line in lines if line.strip()]
        return titles[:3]
    except Exception as e:
        app.logger.exception("Error processing titles")
        return [text]

def process_bullet_points(text):
    try:
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if line:
                line = re.sub(r'[\*\[\]]', '', line)
                line = re.sub(r'^[\d\.\-\s]+', '', line)
                cleaned_lines.append('- ' + line.strip())
        return '\n'.join(cleaned_lines)
    except Exception as e:
        app.logger.exception("Error processing bullet points")
        return text

def generate_slide_titles(content, language="en", desired_count=5):
    try:
        # Explicitly specify the language in the prompt
        lang_name = "Hindi" if language == "hi" else "Telugu" if language == "te" else language.capitalize()
        
        # Request more diverse titles based on the desired count
        prompt = f"""Generate exactly {desired_count} unique and interesting slide titles for a presentation on '{content}' in {lang_name}.
        The titles should:
        - Be concise (3-6 words each)
        - Cover different aspects of the topic
        - Not be repetitive or use numbering
        - Be engaging and informative
        
        Format as a simple list with one title per line, no preamble or extra formatting."""
        
        app.logger.debug(f"Generating {desired_count} titles with prompt: {prompt}")
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        titles_text = response.text.strip()
        app.logger.debug(f"Raw titles response: {titles_text}")
        
        titles = process_titles(titles_text)
        
        # Ensure we have the requested number of unique titles
        seen_titles = set()
        unique_titles = []
        for title in titles:
            if title.lower() not in seen_titles:
                unique_titles.append(title)
                seen_titles.add(title.lower())
        
        # Add more diverse titles if needed
        while len(unique_titles) < desired_count:
            index = len(unique_titles) + 1
            new_aspects = ["Applications", "Benefits", "Challenges", "Future", "Implementation", 
                           "History", "Case Studies", "Best Practices", "Technologies", "Impact", 
                           "Strategies", "Examples", "Analysis", "Methods", "Innovations"]
            aspect = new_aspects[(index - 1) % len(new_aspects)]
            new_title = f"{content} {aspect}"
            if new_title.lower() not in seen_titles:
                unique_titles.append(new_title)
                seen_titles.add(new_title.lower())
        
        return unique_titles[:desired_count]
    except Exception as e:
        app.logger.exception(f"Failed to generate titles: {str(e)}")
        # If title generation fails, create diverse fallback titles
        fallback_titles = []
        aspects = ["Overview", "Introduction", "Applications", "Benefits", "Challenges", 
                  "Future Trends", "Implementation", "Case Studies", "Best Practices", "Impact"]
        for i in range(min(desired_count, len(aspects))):
            fallback_titles.append(f"{content} {aspects[i]}")
        return fallback_titles[:desired_count]

def generate_slide_content(slide_title, has_image=True, summarize=False, language="en"):
    try:
        # Explicitly specify the language in the prompt
        lang_name = "Hindi" if language == "hi" else "Telugu" if language == "te" else language.capitalize()
        if summarize:
            prompt = f"Summarize content for '{slide_title}' into two concise paragraphs (max 30 words each) in {lang_name}, no preamble or labels."
            max_tokens = 100
        elif has_image:
            prompt = f"Generate two concise paragraphs (max 50 words each) for '{slide_title}' in {lang_name}, no preamble or labels."
            max_tokens = 150
        else:
            prompt = f"Generate six concise bullet points (max 25 words each) for '{slide_title}' in {lang_name}. Use '-' as bullet marker, no numbering."
            max_tokens = 300
        app.logger.debug(f"Generating content with prompt: {prompt}")
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        content = response.text.strip()
        app.logger.debug(f"Raw content response: {content}")
        if not has_image and not summarize:
            content = process_bullet_points(content)
        return content
    except Exception as e:
        app.logger.exception(f"Failed to generate content for '{slide_title}': {str(e)}")
        return f"Content generation failed: {str(e)}"

# Improved function to generate_image without any placeholder fallbacks
def generate_image(prompt, language="en"):
    try:
        stability_api_key = os.environ.get('STABILITY_API_KEY')
        
        # If no API key, skip image generation
        if not stability_api_key:
            app.logger.error("STABILITY_API_KEY is not set in the environment!")
            return None
        
        # Call the Stability API
        api_url = "https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image"
        headers = {"Authorization": f"Bearer {stability_api_key}", "Content-Type": "application/json"}
        
        # Customize prompt based on language
        style_prompt = "professional high-quality illustration"
        if language == "hi":
            style_prompt += ", styled for Hindi audience"
        elif language == "te":
            style_prompt += ", styled for Telugu audience"
        else:
            style_prompt += f", styled for {language} audience"
        
        # Set API parameters
        payload = {
            "text_prompts": [{"text": f"{prompt}, {style_prompt}"}],
            "cfg_scale": 7,
            "height": 1024,
            "width": 1024,
            "samples": 1,
            "steps": 30,
        }
        
        # Make the API call
        app.logger.debug(f"Requesting image with data: {payload}")
        response = requests.post(api_url, headers=headers, json=payload)
        response.raise_for_status()
        
        # Parse the response
        data = response.json()
        
        if "artifacts" in data and len(data["artifacts"]) > 0:
            # Get the base64-encoded image
            image_b64 = data["artifacts"][0]["base64"]
            image_data = base64.b64decode(image_b64)
            
            # Convert to BytesIO for PPTX to use
            image_stream = BytesIO(image_data)
            image_stream.seek(0)
            
            app.logger.debug(f"Image generated successfully for prompt: {prompt}")
            return image_stream
        else:
            app.logger.warning(f"No image data in response for prompt: {prompt}")
            return None
            
    except Exception as e:
        app.logger.exception(f"Error generating image: {str(e)}")
        return None

def generate_chart(slide, csv_file, chart_type="bar", language="en"):
    try:
        df = pd.read_csv(csv_file)
        chart_data = CategoryChartData()
        chart_data.categories = df.iloc[:, 0].tolist()
        chart_data.add_series('Data', df.iloc[:, 1].tolist())
        chart_type_enum = CHART_TYPES.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart = slide.shapes.add_chart(
            chart_type_enum,
            Inches(5.5), Inches(1.2),
            Inches(4), Inches(3),
            chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "डेटा अवलोकन" if language == "hi" else "డేటా అవలోకనం" if language == "te" else "Data Overview"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else "Arial"
    except Exception as e:
        app.logger.exception("Error generating chart")

# Add a function to create custom shapes and design elements
def add_design_elements(slide, theme, variant_info):
    try:
        # Add accent bar or graphic based on layout style
        layout_style = variant_info.get("layout", "clean")
        
        if layout_style in ["clean", "professional", "structured"]:
            # Add a subtle accent bar on the left
            left = Inches(0.1)
            top = Inches(0.5)
            width = Inches(0.2)
            height = Inches(6.5)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["accent_color"]
            shape.line.fill.background()
            
        elif layout_style in ["dynamic", "creative", "asymmetric"]:
            # Add diagonal accent element
            left = Inches(7.5)
            top = Inches(5.0)
            width = Inches(3.0)
            height = Inches(2.0)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["accent_color"]
            shape.fill.transparency = 0.7
            shape.line.fill.background()
            
        elif layout_style in ["minimal", "simple", "spacious"]:
            # Add subtle footer line
            left = Inches(0.5)
            top = Inches(7.0)
            width = Inches(9.0)
            height = Inches(0.03)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme["accent_color"]
            shape.line.fill.background()
    except Exception as e:
        # If any error occurs, just continue without the design element
        app.logger.warning(f"Could not add design elements: {str(e)}")

# Enhanced styling function with more distinctive visual elements
def apply_variant_styling(slide, theme, variant_info, is_title_slide=False):
    """Apply variant-specific styling to slide elements with more professional design touches"""
    try:
        layout_style = variant_info.get("layout", "clean")
        bullet_style = variant_info.get("bullet_style", "square")
        
        # PROFESSIONAL VARIANT - Corporate, structured, elegant
        if bullet_style == "square" or bullet_style == "triangle" or bullet_style == "none" or bullet_style == "arrow":
            # Professional style - sophisticated elements
            
            # Add a stylish accent bar
            if not is_title_slide:
                # Add subtle vertical accent bar
                left_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.2),
                    Inches(0.7),
                    Inches(0.08),
                    Inches(6.5)
                )
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = theme["accent_color"]
                left_bar.line.fill.background()
                
                # Add a thin bottom border line
                bottom_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.2),
                    Inches(7.0),
                    Inches(12.9),
                    Inches(0.03)
                )
                bottom_line.fill.solid()
                bottom_line.fill.fore_color.rgb = theme["accent_color"]
                bottom_line.fill.transparency = 0.3
                bottom_line.line.fill.background()
                
                # Add subtle corner accent
                corner_accent = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_TRIANGLE,
                    Inches(12.5),
                    Inches(0),
                    Inches(0.8),
                    Inches(0.8)
                )
                corner_accent.fill.solid()
                corner_accent.fill.fore_color.rgb = theme["accent_color"]
                corner_accent.fill.transparency = 0.7
                corner_accent.line.fill.background()
            
            # For title slide, add more dramatic elements
            if is_title_slide:
                # Add an underline for the title
                title_underline = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.0),
                    Inches(2.5),
                    Inches(6.0),
                    Inches(0.05)
                )
                title_underline.fill.solid()
                title_underline.fill.fore_color.rgb = theme["accent_color"]
                title_underline.line.fill.background()
                
                # Add decorative corner elements
                for i in range(3):
                    size = Inches(0.7 - i*0.2)
                    corner_box = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(13.33 - size),
                        Inches(7.5 - size),
                        size,
                        size
                    )
                    corner_box.fill.solid()
                    corner_box.fill.fore_color.rgb = theme["accent_color"]
                    corner_box.fill.transparency = 0.3 + (i * 0.2)
                    corner_box.line.fill.background()
        
        # CREATIVE VARIANT - Dynamic, artistic, expressive
        elif bullet_style == "circle" or bullet_style == "star" or bullet_style == "checkmark" or bullet_style == "diamond":
            # Creative style - artistic elements
            
            if not is_title_slide:
                # Add decorative circular elements
                for i in range(4):
                    size = Inches(0.4)
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(12.5 - i*0.4),
                        Inches(0.3 + i*0.3),
                        size,
                        size
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = theme["accent_color"]
                    circle.fill.transparency = 0.2 + (i * 0.1)
                    circle.line.fill.background()
                
                # Add swoosh design element
                swoosh = slide.shapes.add_shape(
                    MSO_SHAPE.CURVED_RIGHT_ARROW,
                    Inches(10.0),
                    Inches(6.5),
                    Inches(3.0),
                    Inches(0.8)
                )
                swoosh.fill.solid()
                swoosh.fill.fore_color.rgb = theme["accent_color"]
                swoosh.fill.transparency = 0.6
                swoosh.line.fill.background()
                
                # Add artistic corner element
                for i in range(3):
                    diamond = slide.shapes.add_shape(
                        MSO_SHAPE.DIAMOND,
                        Inches(0.2 + i*0.3),
                        Inches(6.8 - i*0.3),
                        Inches(0.3),
                        Inches(0.3)
                    )
                    diamond.fill.solid()
                    diamond.fill.fore_color.rgb = theme["accent_color"]
                    diamond.fill.transparency = 0.4
                    diamond.line.fill.background()
            
            # For title slide, add artistic elements
            if is_title_slide:
                # Add curved design element
                arc = slide.shapes.add_shape(
                    MSO_SHAPE.ARC,
                    Inches(9.0),
                    Inches(1.0),
                    Inches(4.0),
                    Inches(4.0)
                )
                arc.fill.solid()
                arc.fill.fore_color.rgb = theme["accent_color"]
                arc.fill.transparency = 0.7
                arc.line.fill.background()
                
                # Add overlapping circles
                for i in range(3):
                    size = Inches(1.5)
                    pos_x = Inches(11.0 - i*0.8)
                    pos_y = Inches(5.5 + i*0.4)
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        pos_x, pos_y, size, size
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = theme["accent_color"]
                    circle.fill.transparency = 0.6 + (i * 0.1)
                    circle.line.fill.background()
        
        # MINIMAL VARIANT - Simple, clean, spacious
        elif bullet_style == "dash" or bullet_style == "hyphen" or bullet_style == "dot":
            # Minimal style - subtle, clean elements
            
            if not is_title_slide:
                # Add subtle top border
                top_border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5),
                    Inches(0.3),
                    Inches(12.33),
                    Inches(0.02)
                )
                top_border.fill.solid()
                top_border.fill.fore_color.rgb = theme["accent_color"]
                top_border.fill.transparency = 0.5
                top_border.line.fill.background()
                
                # Add minimal corner marker
                corner = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0),
                    Inches(0),
                    Inches(0.5),
                    Inches(0.06)
                )
                corner.fill.solid()
                corner.fill.fore_color.rgb = theme["accent_color"]
                corner.fill.transparency = 0.3
                corner.line.fill.background()
                
                # Add vertical accent
                vert_accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0),
                    Inches(0),
                    Inches(0.06),
                    Inches(0.5)
                )
                vert_accent.fill.solid()
                vert_accent.fill.fore_color.rgb = theme["accent_color"]
                vert_accent.fill.transparency = 0.3
                vert_accent.line.fill.background()
            
            # For title slide, add minimal style elements
            if is_title_slide:
                # Add thin line above title
                title_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.0),
                    Inches(1.8),
                    Inches(4.0),
                    Inches(0.03)
                )
                title_line.fill.solid()
                title_line.fill.fore_color.rgb = theme["accent_color"]
                title_line.line.fill.background()
                
                # Add minimal dot decoration
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.5),
                    Inches(1.7),
                    Inches(0.15),
                    Inches(0.15)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = theme["accent_color"]
                dot.line.fill.background()
    
    except Exception as e:
        app.logger.warning(f"Error applying variant styling: {str(e)}")

# Add a function to improve bullet point styling
def apply_bullet_styling(paragraph, bullet_style, theme_color):
    """Apply custom bullet styling based on the variant"""
    try:
        if bullet_style == "square":
            # Professional square bullets
            paragraph.bullet.character = '■'
        elif bullet_style == "circle":
            # Creative circle bullets
            paragraph.bullet.character = '●'
        elif bullet_style == "dash":
            # Minimal dash bullets
            paragraph.bullet.character = '—'
        elif bullet_style == "arrow":
            # Arrow bullets
            paragraph.bullet.character = '➔'
        elif bullet_style == "star":
            # Star bullets
            paragraph.bullet.character = '★'
        elif bullet_style == "checkmark":
            # Checkmark bullets
            paragraph.bullet.character = '✓'
        elif bullet_style == "diamond":
            # Diamond bullets
            paragraph.bullet.character = '◆'
        elif bullet_style == "triangle":
            # Triangle bullets
            paragraph.bullet.character = '▶'
        elif bullet_style == "dot":
            # Dot bullets
            paragraph.bullet.character = '•'
        else:
            # Default bullet
            paragraph.bullet.character = '•'
            
        # Set bullet color to match theme accent
        paragraph.bullet.font.color.rgb = theme_color
    except:
        # If custom bullets fail, use default bullet
        paragraph.bullet = True

# Update the create_presentation function to use these enhancements
def create_presentation(topic, text_file=None, csv_file=None, theme="corporate", variant="professional", language="en", include_images=True, summarize=False, chart_type="bar", export_format="pptx", slide_count=5):
    try:
        prs = Presentation()
        selected_theme = THEMES.get(theme, THEMES["corporate"])
        variant_info = selected_theme["variants"].get(variant, selected_theme["variants"]["professional"])
        
        # Set slide size to widescreen (16:9) for a more modern look
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        content = topic
        if text_file:
            content = text_file.read().decode('utf-8')
        app.logger.debug(f"Content: {content[:100]}...")

        # Get slide titles with improved function
        desired_content_slides = min(int(slide_count), 10) - 1  # -1 for title slide
        
        # Generate slide titles
        slide_titles = generate_slide_titles(content, language, desired_content_slides + 1)
        
        # Log all slide titles for debugging
        app.logger.debug(f"Generated slide titles: {slide_titles}")
        
        # Ensure we have enough titles
        while len(slide_titles) < desired_content_slides + 1:
            # Add generic but topic-related titles if we don't have enough
            default_aspects = ["Overview", "Applications", "Benefits", "Challenges", "Future Trends", 
                            "Implementation", "Case Studies", "Best Practices", "Impact", "Technologies"]
            index = len(slide_titles)
            new_title = f"{topic} {default_aspects[index % len(default_aspects)]}"
            slide_titles.append(new_title)
            app.logger.debug(f"Added fallback title: {new_title}")
        
        # Title slide - no image on first slide
        # Use blank layout (layout 6) instead of title layout to avoid placeholders
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Using blank layout
        apply_gradient(title_slide, selected_theme["gradient_start"], selected_theme["gradient_end"])
        
        # Add design elements
        add_design_elements(title_slide, selected_theme, variant_info)
        apply_variant_styling(title_slide, selected_theme, variant_info, is_title_slide=True)
        
        # Manually add title to blank slide
        title_shape = title_slide.shapes.add_textbox(
            Inches(1.0),
            Inches(2.5),
            Inches(10.0),
            Inches(1.5)
        )
        title_tf = title_shape.text_frame
        title_tf.text = slide_titles[0]
        title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
        title_tf.paragraphs[0].font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else variant_info["font_name"]
        title_tf.paragraphs[0].font.color.rgb = selected_theme["title_color"]
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center the title
        
        if selected_theme["shadow"]:
            add_text_shadow(title_tf)
            
        # Add subtitle
        subtitle_shape = title_slide.shapes.add_textbox(
            Inches(2.0),
            Inches(4.0),
            Inches(8.0),
            Inches(1.0)
        )
        subtitle_tf = subtitle_shape.text_frame
        subtitle_tf.text = "Powered by AI"
        subtitle_tf.paragraphs[0].font.size = Pt(24)
        subtitle_tf.paragraphs[0].font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else variant_info["font_name"]
        subtitle_tf.paragraphs[0].font.color.rgb = selected_theme["text_color"]
        subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Content slides with improved layout and design
        for i, title in enumerate(slide_titles):
            # Create slide with consistent layout
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Using blank layout
            apply_gradient(slide, selected_theme["gradient_start"], selected_theme["gradient_end"])
            
            # Add design elements
            add_design_elements(slide, selected_theme, variant_info)
            apply_variant_styling(slide, selected_theme, variant_info)
            
            # Add title manually
            title_shape = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.5),
                Inches(11.0),
                Inches(0.8)
            )
            title_tf = title_shape.text_frame
            title_tf.text = title
            title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
            title_tf.paragraphs[0].font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else variant_info["font_name"]
            title_tf.paragraphs[0].font.color.rgb = selected_theme["title_color"]
            title_tf.paragraphs[0].font.bold = True
            
            # Generate content for the slide
            content_text = generate_slide_content(title, include_images, language)
            app.logger.debug(f"Generated content for slide {i+1}: {title}")
            
            # Position text on left side if there's an image
            if include_images and (i % 2 == 0):  # Every other slide gets an image
                # Narrower text box on the left for slides with images
                textbox = slide.shapes.add_textbox(
                    Inches(0.5),       # Left position
                    Inches(1.5),       # Top position
                    Inches(5.75),      # Width (about half the slide)
                    Inches(5.0)        # Height
                )
            else:
                # Full-width text box for text-only slides
                textbox = slide.shapes.add_textbox(
                    Inches(0.7),       # Left position
                    Inches(1.5),       # Top position
                    Inches(11.0),      # Width (full slide width)
                    Inches(5.0)        # Height
                )
            
            # Configure text frame
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.text = content_text
            
            # Format content based on type (paragraphs or bullet points)
            if content_text.startswith('-'):
                # For bullet points
                tf.text = ""  # Clear the text frame
                lines = content_text.split('\n')
                for line_idx, line in enumerate(lines):
                    p = tf.add_paragraph()
                    p.text = line.lstrip('- ')
                    p.level = 0
                    p.font.size = CONTENT_FONT_SIZE
                    p.font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else variant_info["font_name"]
                    p.font.color.rgb = selected_theme["text_color"]
                    
                    # Add enhanced bullet formatting based on variant style
                    bullet_style = variant_info.get("bullet_style", "square")
                    apply_bullet_styling(p, bullet_style, selected_theme["accent_color"])
            else:
                # For regular paragraphs
                for paragraph in tf.paragraphs:
                    paragraph.font.size = CONTENT_FONT_SIZE
                    paragraph.font.name = "Mangal" if language == "hi" else "Gautami" if language == "te" else variant_info["font_name"]
                    paragraph.font.color.rgb = selected_theme["text_color"]
                    paragraph.space_after = Pt(12)
            
            if selected_theme["shadow"]:
                add_text_shadow(tf)
            
            # Add image if needed
            if include_images and (i % 2 == 0):  # Every other slide gets an image
                image_stream = generate_image(f"{title} related to {content_text}", language)
                if image_stream:
                    left_img = Inches(7.0)      # Position to the right
                    top_img = Inches(1.5)       # Same top position as text
                    width_img = Inches(5.5)     # Width of image
                    slide.shapes.add_picture(image_stream, left_img, top_img, width=width_img)
                    app.logger.debug(f"Image added to slide: {title}")

        # Add footer to all slides
        for slide in prs.slides:
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.3))
            tf = footer.text_frame
            tf.text = f"{topic} | {pd.Timestamp.now().strftime('%Y-%m-%d')}"
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = selected_theme["text_color"]
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Save the presentation
        output_dir = "generated_ppt"
        os.makedirs(output_dir, exist_ok=True)
        pptx_filepath = os.path.join(output_dir, f"{topic or 'presentation'}_presentation.pptx")
        prs.save(pptx_filepath)
        app.logger.debug(f"Saved to {pptx_filepath}")
        
        return pptx_filepath
    except Exception as e:
        app.logger.exception("Error in create_presentation")
        raise e

# Add this function to check if a slide has the expected placeholders
def ensure_slide_has_title(slide):
    """Make sure the slide has a title placeholder, adding one if missing"""
    has_title = False
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # 1 is title
            has_title = True
            break
    
    if not has_title:
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0)
        )
        title_shape.text_frame.text = "Slide Title"
        return title_shape
    
    return slide.shapes.title

@app.route('/generate', methods=['POST'])
def generate():
    topic = request.form.get('topic')
    text_file = request.files.get('textFile')
    csv_file = request.files.get('csvFile')
    theme = request.form.get('theme', 'corporate')
    variant = request.form.get('variant', 'professional')
    language = request.form.get('language', 'en')
    include_images = request.form.get('includeImages', 'true') == 'true'
    summarize = request.form.get('summarize', 'false') == 'true'
    chart_type = request.form.get('chartType', 'bar')
    export_format = request.form.get('exportFormat', 'pptx')
    slide_count = request.form.get('slideCount', '5')  # Default to 5 slides
    
    # Validate slide count
    try:
        slide_count = int(slide_count)
        if slide_count < 3:
            slide_count = 3  # Minimum 3 slides
        elif slide_count > 10:
            slide_count = 10  # Maximum 10 slides
    except ValueError:
        slide_count = 5  # Default if invalid

    if not topic and not text_file:
        return jsonify({"error": "Topic or text file required"}), 400

    try:
        app.logger.debug(f"Generating presentation with {slide_count} slides")
        output_filepath = create_presentation(
            topic=topic,
            text_file=text_file,
            csv_file=csv_file,
            theme=theme,
            variant=variant,
            language=language,
            include_images=include_images,
            summarize=summarize,
            chart_type=chart_type,
            export_format=export_format,
            slide_count=slide_count
        )
        
        download_name = f"{topic or 'presentation'}_presentation"
        
        # Handle different export formats
        if export_format == 'pdf':
            try:
                pdf_filepath = output_filepath.replace('.pptx', '.pdf')
                
                # Create a PDF document
                pdf = FPDF()
                pdf.set_auto_page_break(auto=True, margin=15)
                
                # Load the presentation to extract content
                prs = Presentation(output_filepath)
                
                # First page - title
                pdf.add_page()
                pdf.set_font("Arial", 'B', 24)
                pdf.cell(0, 20, txt=topic or "AI-Generated Presentation", ln=True, align='C')
                pdf.set_font("Arial", 'I', 14)
                pdf.cell(0, 10, txt=f"Created on {pd.Timestamp.now().strftime('%Y-%m-%d')}", ln=True, align='C')
                
                # Add slide content to PDF
                for i, slide in enumerate(prs.slides):
                    pdf.add_page()
                    
                    # Add slide number
                    pdf.set_font("Arial", 'I', 10)
                    pdf.cell(0, 10, txt=f"Slide {i+1}", ln=True, align='R')
                    
                    # Add slide title if present
                    if len(slide.shapes.title.text_frame.text) > 0:
                        pdf.set_font("Arial", 'B', 16)
                        pdf.cell(0, 15, txt=slide.shapes.title.text_frame.text, ln=True)
                    
                    # Add slide content - extract text from shapes
                    pdf.set_font("Arial", '', 12)
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                            if shape.text_frame.text and shape != slide.shapes.title:
                                # Process each paragraph
                                for paragraph in shape.text_frame.paragraphs:
                                    text = paragraph.text.strip()
                                    if text:
                                        pdf.multi_cell(0, 8, txt=text)
                                        pdf.ln(4)
                
                # Save the PDF
                pdf.output(pdf_filepath)
                
                app.logger.debug(f"Created PDF document: {pdf_filepath}")
                
                # If the PDF was created successfully, send it
                if os.path.exists(pdf_filepath):
                    return send_file(pdf_filepath, as_attachment=True, download_name=f"{download_name}.pdf")
                else:
                    # Fall back to PPTX if PDF creation failed
                    app.logger.warning("PDF file was not created, falling back to PPTX")
                    return send_file(output_filepath, as_attachment=True, download_name=f"{download_name}.pptx")
            except Exception as e:
                app.logger.exception("PDF conversion failed, falling back to PPTX")
                # If PDF conversion fails, fall back to PPTX
                return send_file(output_filepath, as_attachment=True, download_name=f"{download_name}.pptx")
        else:  # Default to PPTX
            return send_file(output_filepath, as_attachment=True, download_name=f"{download_name}.pptx")
    except Exception as e:
        app.logger.exception("Error generating presentation")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(port=5000, debug=True)